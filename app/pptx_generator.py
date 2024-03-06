#!/usr/bin/env python3

from pptx import Presentation
from datetime import datetime, timedelta
import os
import sys
import yaml
from importlib import resources

# print(f"Current Working Directory: {os.getcwd()}")
script_path = os.path.abspath(__file__)
script_dir = os.path.split(script_path)[0]
# print(script_dir)
# print(f"Python Interpreter: {sys.executable}")
# pptx_file = importlib.resources.resource_filename('app', 'data/Remaja_template.pptx')
# with resources.files('app.data').joinpath('Remaja_template.pptx') as pptx_path:
#     prs = Presentation(pptx_path)
#     pptx_file = str(pptx_path)
pptx_path = os.path.join(script_dir, "data/Remaja_template.pptx")

with open(os.path.join(script_dir, "data/data.yml"), 'r') as file:
    data = yaml.safe_load(file)

def get_next_weekday(startdate, weekday):
    """
    @startdate: given date, in format '2013-05-25'
    @weekday: week day as a integer, between 0 (Monday) to 6 (Sunday)
    """
    d = datetime.strptime(startdate, '%d %B %Y')
    t = timedelta((7 + weekday - d.weekday()) % 7)
    return (d + t).strftime('%d %B %Y')

prs = Presentation(pptx_path)
output_presentation = './presentation.pptx'
# ------------------- Song ID -------------------------------
song_1_id = '003'
song_2_id = '001'
song_3_id = '004'

def id_to_song():
    global song_1
    global song_2 
    global song_3
    print(data)
    song_1 = data[song_1_id]
    song_2 = data[song_2_id]
    song_3 = data[song_3_id]

def generate_verse_number(verse_num): 
    if isinstance(verse_num, int):
        return str(verse_num)
    elif isinstance(verse_num, str):
        return verse_num
    else: 
        return 0
def add_welcome_slide():
    # print("Welcome Slide")
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    # for shape in slide.placeholders:
    #  print('%d %s' % (shape.placeholder_format.idx, shape.name))
    date_text = slide.placeholders[10]
    church_name = slide.placeholders[11]
    title.text = "Welcome to Remaja"
    # date_text.text = datetime.strptime(saturday,'%d %B %Y')
    date_text.text = get_next_weekday(datetime.today().strftime('%d %B %Y'), 5)
    church_name.text = "Reformed Evangelical Church Singapore"

def add_first_song_title_slide():
    # print("First Slide")
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    # for shape in slide.placeholders:
    #  print('%d %s' % (shape.placeholder_format.idx, shape.name))
    title.text = song_1['title']
    number_text = slide.placeholders[10]
    author_text = slide.placeholders[11]
    author_text.text = song_1['author']
    number_text.text = "1"

def add_second_song_title_slide():
    # print("Second Slide")
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    title = slide.shapes.title
    # for shape in slide.placeholders:
    #  print('%d %s' % (shape.placeholder_format.idx, shape.name))
    title.text = song_2['title']
    number_text = slide.placeholders[10]
    author_text = slide.placeholders[11]
    author_text.text = song_2['author']
    number_text.text = "2"

def add_third_song_title_slide():
    # print("Third Slide")
    slide = prs.slides.add_slide(prs.slide_layouts[3])
    title = slide.shapes.title
    # for shape in slide.placeholders:
    #  print('%d %s' % (shape.placeholder_format.idx, shape.name))
    title.text = song_3['title']
    number_text = slide.placeholders[10]
    author_text = slide.placeholders[11]
    author_text.text = song_3['author']
    number_text.text = "3"
    
def add_first_song_lyrics_slide(verse_number):
    # print("First Lyric Slide")
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    title = slide.shapes.title
    # for shape in slide.placeholders:
    #  print('%d %s' % (shape.placeholder_format.idx, shape.name))
    title.text = song_1['title']
    author_text = slide.placeholders[11]
    number_text = slide.placeholders[13]
    lyrics_text = slide.placeholders[12].text_frame
    author_text.text = song_1['author']
    # print(verse_number)
    lyrics_text.text = song_1['lyrics'][verse_number][0]
    number_text.text = generate_verse_number(song_1['verse_number'][verse_number])
    for i in range(1, len(song_1['lyrics'][verse_number])):
        p = lyrics_text.add_paragraph()
        p.text = song_1['lyrics'][verse_number][i]

def add_second_song_lyrics_slide(verse_number):
    # print("First Lyric Slide")
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    # for shape in slide.placeholders:
    #  print('%d %s' % (shape.placeholder_format.idx, shape.name))
    title.text = song_2['title']
    author_text = slide.placeholders[11]
    number_text = slide.placeholders[13]
    lyrics_text = slide.placeholders[12].text_frame
    author_text.text = song_2['author']
    number_text.text = generate_verse_number(song_2['verse_number'][verse_number])
    lyrics_text.text = song_2['lyrics'][verse_number][0]
    for i in range(1, len(song_2['lyrics'][verse_number])):
        p = lyrics_text.add_paragraph()
        p.text = song_2['lyrics'][verse_number][i]

def add_third_song_lyrics_slide(verse_number):
    # print("First Lyric Slide")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.title
    # for shape in slide.placeholders:
    #  print('%d %s' % (shape.placeholder_format.idx, shape.name))
    title.text = song_3['title']
    author_text = slide.placeholders[11]
    number_text = slide.placeholders[13]
    lyrics_text = slide.placeholders[12].text_frame
    author_text.text = song_3['author']
    number_text.text = generate_verse_number(song_3['verse_number'][verse_number])
    lyrics_text.text = song_3['lyrics'][verse_number][0]
    for i in range(1, len(song_3['lyrics'][verse_number])):
        p = lyrics_text.add_paragraph()
        p.text = song_3['lyrics'][verse_number][i]    

def add_black_slide():
    prs.slides.add_slide(prs.slide_layouts[7])

def add_announcements():
    prs.slides.add_slide(prs.slide_layouts[8])
    prs.slides.add_slide(prs.slide_layouts[9])
    birthday = prs.slides.add_slide(prs.slide_layouts[10])
    birthday.shapes.title.text = "Birthday Song"
    prs.slides.add_slide(prs.slide_layouts[11])

def generate_first_song():
    add_first_song_title_slide()
    for i in range(0, len(song_1['verse_number'])):
        add_first_song_lyrics_slide(i)
    add_black_slide()

def generate_second_song():
    add_second_song_title_slide()
    for i in range(0, len(song_2['verse_number'])):
        add_second_song_lyrics_slide(i)
    add_black_slide()

def generate_third_song():
    add_third_song_title_slide()
    for i in range(0, len(song_3['verse_number'])):
        add_third_song_lyrics_slide(i)
    add_black_slide()

def song_id():
    global song_1_id
    global song_2_id
    global song_3_id
    global Fail
    Fail = 0
    if data.get(sys.argv[1]):
        song_1_id = sys.argv[1]
        Fail += 1
    else: 
        song_1_id = sys.argv[1]
        print(f"{song_1_id} is not a valid song id" )
    if data.get(sys.argv[2]):
        song_2_id = sys.argv[2]
        Fail += 1
    else: 
        song_2_id = sys.argv[2]
        print(f"{song_2_id} is not a valid song id" )
    if data.get(sys.argv[3]):
        song_3_id = sys.argv[3]
        Fail += 1
    else: 
        song_3_id = sys.argv[3]
        print(f"{song_3_id} is not a valid song id" )
    if Fail != 3: 
        sys.exit()

#----------------------------- MAIN FUNCTION -----------------------------------------------------------------

def main():
    if len(sys.argv) == 4:
        # song_1_id = sys.argv[1]
        # song_2_id = sys.argv[2]
        # song_3_id = sys.argv[3]
        song_id()
        id_to_song()
        add_welcome_slide()
        generate_first_song()
        generate_second_song()
        generate_third_song() 
        add_announcements() 
        prs.save(output_presentation)
        print("Done!")
    elif len(sys.argv) > 4: 
        # song_1_id = sys.argv[1]
        # song_2_id = sys.argv[2]
        # song_3_id = sys.argv[3]
        song_id()
        if sys.argv[4].endswith(('.pptx')):
            output_presentation = sys.argv[4]
        else:
            output_presentation = sys.argv[4] + ".pptx"
        id_to_song()
        add_welcome_slide()
        generate_first_song()
        generate_second_song()
        generate_third_song() 
        add_announcements() 
        prs.save(output_presentation)
        print("Done!")
       
    else:
        print("Not enough arguments provided. Requires at least 3 values")
        print(sys.argv)

def search():
    if len(sys.argv) == 3: 
        found = 0
        if sys.argv[2] == "-a":
            for key in data:
                if sys.argv[1].lower().replace(" ", "") == data[key]['author'].lower().replace(" ", ""):
                    print(key)
                    found = 1
            if found != 1: 
                print(f"There is no '{sys.argv[1]}' in our song list.")
        elif sys.argv[2] == "-t":
            for key in data:
                if sys.argv[1].lower().replace(" ", "") == data[key]['title'].lower().replace(" ", ""):
                    print(key)
                    found = 1
            if found != 1: 
                print(f"There is no '{sys.argv[1]}' in our song list.")
        elif sys.argv[2] == "-k":
            if sys.argv[2].isdigit():
                for key in data:
                    if sys.argv[1] == data[key]['kri_number']:
                        print(key)
                        found = 1
                if found != 1: 
                    print(f"There is no '{sys.argv[1]}' in our song list.")
            else:
                print("That is not a valid KRI number")
    else: 
        print('pptx-generator-search requires 2 values. \nThe first is the author, title, or KRI number of the song surrounded in double quotes ("), \nNext is the flag. -a for author, -t for title, and -k for kri number\nExample: pptx-generator-search "Grace Alone" -t')
if __name__ == "__main__":
    main()
